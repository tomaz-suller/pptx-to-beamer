from dataclasses import dataclass, field
from hashlib import sha1
from pathlib import Path
from typing import ClassVar, Union

import pptx

from text_classifiers import is_accent, is_strong, is_title


@dataclass
class Renderable:
    save_dirpath: Path = field(init=False, repr=False)

    def to_markdown(self) -> str:
        raise NotImplementedError


class Parser:
    @staticmethod
    def from_pptx(obj):
        raise NotImplementedError


@dataclass
class BulletPoint:
    level: int
    text: str


@dataclass
class BulletPoints(Renderable):
    bullets: list[BulletPoint]

    def to_markdown(self) -> str:
        result = [
            bullet_point.level * '\t'
            + '* ' + bullet_point.text
            for bullet_point in self.bullets
        ]
        return '\n'.join(result)


@dataclass
class Paragraphs(Renderable):
    paragraphs: list[str]

    def to_markdown(self) -> str:
        return '\n'.join(self.paragraphs)


@dataclass
class Title(Parser, Renderable):
    text: str

    def to_markdown(self) -> str:
        return '## ' + self.text

    @staticmethod
    def from_pptx(obj) -> 'Title':
        return Title(obj.text.strip())


@dataclass
class Text(Parser, Renderable):
    is_bullet: bool
    content: Union[Paragraphs, BulletPoints]

    def to_markdown(self) -> str:
        self.content.to_markdown()

    @staticmethod
    def from_pptx(obj) -> 'Text':
        paragraphs = obj.text_frame.paragraphs
        paragraphs_text = [
            Text._parse_runs(paragraph.runs)
            for paragraph in paragraphs
        ]
        paragraphs_level = [
            paragraph.level
            for paragraph in paragraphs
        ]
        if len(set(paragraphs_level)) > 1 or len(paragraphs) > 2:
            content = BulletPoints([
                BulletPoint(level, text)
                for level, text in zip(paragraphs_level, paragraphs_text)
            ])
        else:
            content = Paragraphs(paragraphs_text)

        return content

    @staticmethod
    def _parse_runs(runs: list) -> str:

        result = []
        for run in runs:
            emphasis_markers = (
                '_'
                * (is_accent(run.font) + is_strong(run.font))
            )
            result.append(
                emphasis_markers
                + run.text.strip()
                + emphasis_markers
            )
        return ' '.join(result)


@dataclass
class Image(Parser, Renderable):
    filename: str
    image: bytes = field(repr=False)
    IMAGE_DIR: ClassVar[str] = 'img'

    def to_markdown(self) -> str:
        relative_image_filepath = Path(
            self.IMAGE_DIR,
            self.filename
        )
        image_filepath = self.save_dirpath / relative_image_filepath
        image_filepath.parent.mkdir(parents=True, exist_ok=True)
        with open(image_filepath, 'wb') as f:
            f.write(self.image)
        return f'![]({relative_image_filepath})'

    @staticmethod
    def from_pptx(obj) -> 'Image':
        image_bytes = obj.image.blob
        image_hash = sha1()
        image_hash.update(image_bytes)
        image_filename = (
            Path(obj.image.filename).stem
            + '-'
            + image_hash.hexdigest()
            + '.'+obj.image.ext
        )
        return Image(
            filename=image_filename,
            image=image_bytes,
        )


@dataclass
class Slide(Parser):
    title: Title
    content: list[Renderable] = field(repr=False)

    def save(self, output_path: Path) -> str:
        result = [self.title.to_markdown()]
        for renderable in self.content:
            renderable.save_dirpath = output_path
            result.append(renderable.to_markdown())
        return '\n\n'.join(result)

    @staticmethod
    def from_pptx(obj) -> 'Slide':
        content = []
        title = Title('Title placeholder')
        for shape in obj.shapes:
            if hasattr(shape, 'image'):
                content.append(Image.from_pptx(shape))
            elif is_title(shape):
                title = Title.from_pptx(shape)
            elif hasattr(shape, 'text_frame'):
                content.append(Text.from_pptx(shape))

        return Slide(
            title=title,
            content=content,
        )


@dataclass
class Presentation(Parser):
    slides: list[Slide] = field(repr=False)
    OUTPUT_FILENAME: ClassVar[str] = 'out.md'

    def save(self, output_path: Path) -> None:
        result = []
        for slide in self.slides:
            result.append(slide.save(output_path))
        with open(output_path / self.OUTPUT_FILENAME, 'w', encoding='utf8') as f:
            f.write('\n\n\n'.join(result))

    @staticmethod
    def read(path: Path) -> 'Presentation':
        pptx_presentation = pptx.Presentation(path)
        return Presentation.from_pptx(pptx_presentation)

    @staticmethod
    def from_pptx(obj) -> 'Presentation':
        return Presentation([
            Slide.from_pptx(slide)
            for slide in obj.slides
        ])
